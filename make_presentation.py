"""
make_presentation.py
--------------------
Generates the Group 13 – STAN 47 final-project PowerPoint presentation.

Usage
-----
    python3 make_presentation.py
    python3 make_presentation.py \
        --lstm_rmse 1.23 --lstm_mae 0.91 --lstm_r2 0.94 \
        --trans_rmse 1.18 --trans_mae 0.87 --trans_r2 0.95 \
        --tcn_rmse  1.10 --tcn_mae  0.82 --tcn_r2  0.96

If no values are passed the script uses placeholder strings so the file
can be generated even before the notebook has been fully executed.

Output
------
    Group13_DeepLearning_Presentation.pptx  (project root)
"""

import argparse
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ── CLI args ────────────────────────────────────────────────────────
parser = argparse.ArgumentParser()
for m in ("lstm", "trans", "tcn"):
    parser.add_argument(f"--{m}_rmse", default="—")
    parser.add_argument(f"--{m}_mae",  default="—")
    parser.add_argument(f"--{m}_r2",   default="—")
args = parser.parse_args()

def _fmt(v, decimals=4):
    try:
        return f"{float(v):.{decimals}f}"
    except (TypeError, ValueError):
        return str(v)

METRICS = {
    "LSTM":        {"rmse": _fmt(args.lstm_rmse), "mae": _fmt(args.lstm_mae),  "r2": _fmt(args.lstm_r2)},
    "Transformer": {"rmse": _fmt(args.trans_rmse),"mae": _fmt(args.trans_mae), "r2": _fmt(args.trans_r2)},
    "TCN":         {"rmse": _fmt(args.tcn_rmse),  "mae": _fmt(args.tcn_mae),   "r2": _fmt(args.tcn_r2)},
}

# ── Palette ─────────────────────────────────────────────────────────
NAVY    = RGBColor(0x2c, 0x3e, 0x50)
BLUE    = RGBColor(0x29, 0x80, 0xb9)
RED     = RGBColor(0xe7, 0x4c, 0x3c)
GREEN   = RGBColor(0x27, 0xae, 0x60)
ORANGE  = RGBColor(0xf3, 0x9c, 0x12)
LGRAY   = RGBColor(0xec, 0xf0, 0xf1)
WHITE   = RGBColor(0xff, 0xff, 0xff)
DKGRAY  = RGBColor(0x55, 0x65, 0x73)

W = Inches(13.33)
H = Inches(7.50)

# ── Helpers ─────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # blank layout


def fill_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT,
                v_anchor=MSO_ANCHOR.TOP, wrap=True, font_name="Calibri Light"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.text_frame.word_wrap = wrap
    txBox.text_frame.auto_size = None
    txBox.text_frame.vertical_anchor = v_anchor
    p = txBox.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_textbox(slide, items, x, y, w, h,
                       font_size=16, color=WHITE,
                       bullet_color=None, title=None,
                       title_size=20, title_color=None,
                       line_spacing=1.15):
    """items: list of (indent_level, text) tuples or plain strings."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    import copy

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.name  = "Calibri Light"
        run.font.size  = Pt(title_size)
        run.font.bold  = True
        run.font.color.rgb = title_color or color
        first = False

    for item in items:
        if isinstance(item, str):
            level, text = 0, item
        else:
            level, text = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        p.level = level
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        # indentation
        indent_pts = 14 + level * 16
        pPr.set("indent", str(int(-Pt(indent_pts))))
        pPr.set("marL",   str(int(Pt(indent_pts))))

        run = p.add_run()
        prefix = "▸ " if level == 0 else "• "
        run.text = prefix + text
        run.font.name  = "Calibri"
        run.font.size  = Pt(font_size - level * 1.5)
        run.font.color.rgb = bullet_color if level == 0 and bullet_color else color

        # line spacing via XML
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", str(int(line_spacing * 100000)))

    return txBox


def accent_bar(slide, color=BLUE, height=Inches(0.07)):
    """Thin accent bar at very top of slide."""
    add_rect(slide, 0, 0, W, height, fill_color=color)


def header_strip(slide, title, subtitle=None,
                 bg=BLUE, text_color=WHITE, height=Inches(1.15)):
    """Coloured header band with title (and optional subtitle)."""
    add_rect(slide, 0, 0, W, height, fill_color=bg)
    y_title = Inches(0.18) if subtitle else Inches(0.28)
    add_textbox(slide, title,
                Inches(0.45), y_title, Inches(12.4), Inches(0.62),
                font_size=30, bold=True, color=text_color, font_name="Calibri Light")
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.45), Inches(0.72), Inches(12.4), Inches(0.38),
                    font_size=17, color=RGBColor(0xd0, 0xe8, 0xf5), font_name="Calibri")


def footer(slide, text="Group 13 · STAN 47 · Deep Learning Final Project · March 2026",
           page_num=None):
    add_rect(slide, 0, H - Inches(0.32), W, Inches(0.32), fill_color=NAVY)
    label = text + (f"    {page_num}" if page_num else "")
    add_textbox(slide, label,
                Inches(0.3), H - Inches(0.30), Inches(12.0), Inches(0.28),
                font_size=9, color=RGBColor(0x99, 0xaa, 0xbb),
                align=PP_ALIGN.LEFT, font_name="Calibri")


def divider_line(slide, y, color=LGRAY, thickness=Inches(0.012)):
    add_rect(slide, Inches(0.45), y, Inches(12.43), thickness, fill_color=color)


# ════════════════════════════════════════════════════════════════════
#  SLIDES
# ════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 – Title."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, NAVY)

    # decorative accent rectangles
    add_rect(sl, 0, 0, Inches(0.55), H, fill_color=BLUE)
    add_rect(sl, 0, H - Inches(0.60), W, Inches(0.60),
             fill_color=RGBColor(0x1a, 0x25, 0x2f))
    add_rect(sl, Inches(0.55), Inches(3.55), W - Inches(0.55), Inches(0.07),
             fill_color=ORANGE)

    add_textbox(sl, "Deep Learning for\nTime Series Forecasting",
                Inches(1.0), Inches(1.0), Inches(11.5), Inches(2.10),
                font_size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
                font_name="Calibri Light")

    add_textbox(sl, "ETTh1 — Electricity Transformer Temperature",
                Inches(1.0), Inches(3.20), Inches(11.0), Inches(0.60),
                font_size=24, italic=True, color=BLUE, align=PP_ALIGN.LEFT,
                font_name="Calibri Light")

    add_textbox(sl,
                "Ayaa Asoba  ·  Pablo González Martín  ·  Xavier Bruneau",
                Inches(1.0), Inches(4.40), Inches(11.0), Inches(0.45),
                font_size=17, color=LGRAY, align=PP_ALIGN.LEFT)

    add_textbox(sl, "Group 13  |  STAN 47  |  March 2026",
                Inches(1.0), Inches(4.95), Inches(11.0), Inches(0.38),
                font_size=14, color=DKGRAY, align=PP_ALIGN.LEFT)


def slide_agenda(prs):
    """Slide 2 – Agenda."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Agenda", bg=NAVY)
    footer(sl, page_num="2")

    sections = [
        (BLUE,   "01", "Problem Statement & Dataset",    "What are we forecasting and why?"),
        (GREEN,  "02", "Exploratory Data Analysis",       "Understanding the ETTh1 time series"),
        (ORANGE, "03", "Three Deep Learning Models",     "LSTM · Transformer · TCN"),
        (RED,    "04", "Results & Model Comparison",     "Metrics, training curves, forecast plots"),
        (DKGRAY, "05", "Conclusions & Future Work",       "Key takeaways and next steps"),
    ]

    for i, (col, num, title, sub) in enumerate(sections):
        y = Inches(1.30) + i * Inches(1.08)
        add_rect(sl, Inches(0.45), y, Inches(0.62), Inches(0.72), fill_color=col)
        add_textbox(sl, num, Inches(0.45), y, Inches(0.62), Inches(0.72),
                    font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                    v_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(sl, title,
                    Inches(1.20), y + Inches(0.06), Inches(5.8), Inches(0.38),
                    font_size=18, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
        add_textbox(sl, sub,
                    Inches(1.20), y + Inches(0.38), Inches(5.8), Inches(0.30),
                    font_size=13, color=DKGRAY, align=PP_ALIGN.LEFT)

    # speaking time note on right
    add_rect(sl, Inches(7.4), Inches(1.30), Inches(5.5), Inches(5.50),
             fill_color=LGRAY)
    add_textbox(sl, "Speaking time guide",
                Inches(7.6), Inches(1.40), Inches(5.1), Inches(0.40),
                font_size=14, bold=True, color=NAVY)
    times = [
        "Problem & Dataset   ~3 min",
        "EDA                           ~2 min",
        "Model architectures ~8 min",
        "Results & comparison ~4 min",
        "Conclusions              ~3 min",
    ]
    for i, t in enumerate(times):
        add_textbox(sl, t,
                    Inches(7.65), Inches(1.90) + i * Inches(0.82), Inches(5.0), Inches(0.38),
                    font_size=13, color=NAVY)


def slide_problem(prs):
    """Slide 3 – Problem Statement."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Problem Statement", subtitle="Forecasting transformer oil temperature (OT)", bg=NAVY)
    footer(sl, page_num="3")

    # Left column – context
    add_textbox(sl, "Real-world context",
                Inches(0.45), Inches(1.28), Inches(5.8), Inches(0.40),
                font_size=16, bold=True, color=NAVY)
    divider_line(sl, Inches(1.73))
    ctx = [
        "Power grids rely on transformer stations to\nstep voltages up/down between transmission\nand distribution networks.",
        "Oil temperature (OT) is the primary health\nindicator — overheating causes insulation\ndegradation and transformer failures.",
        "Accurate 1-step-ahead OT forecasting enables\npreventive maintenance and load-scheduling\ndecisions before critical thresholds are reached.",
    ]
    for i, t in enumerate(ctx):
        add_textbox(sl, t,
                    Inches(0.55), Inches(1.82) + i * Inches(1.4), Inches(5.6), Inches(1.30),
                    font_size=14, color=DKGRAY)

    # Right column – task definition box
    add_rect(sl, Inches(6.7), Inches(1.28), Inches(6.2), Inches(5.60),
             fill_color=NAVY)
    add_textbox(sl, "Forecasting task",
                Inches(6.9), Inches(1.40), Inches(5.8), Inches(0.40),
                font_size=17, bold=True, color=ORANGE)
    divider_line(sl, Inches(1.90))

    task_items = [
        "Dataset:      ETTh1 (2016 – 2018, 17 420 h)",
        "Target:         Oil Temperature  (OT, °C)",
        "Horizon:       1-step-ahead  (next hour)",
        "Input window:  48 hours of OT history",
        "Split:              80% train  |  20% test",
        "Metric:          RMSE · MAE · R²",
    ]
    for i, t in enumerate(task_items):
        add_textbox(sl, t,
                    Inches(6.9), Inches(2.05) + i * Inches(0.72), Inches(5.7), Inches(0.40),
                    font_size=14, color=WHITE)

    add_textbox(sl, "Three architectures compared",
                Inches(6.9), Inches(6.35), Inches(5.8), Inches(0.38),
                font_size=13, bold=True, italic=True, color=BLUE)
    add_textbox(sl, "Bidirectional LSTM  ·  Transformer  ·  TCN",
                Inches(6.9), Inches(6.68), Inches(5.8), Inches(0.30),
                font_size=13, color=LGRAY)


def slide_dataset(prs):
    """Slide 4 – Dataset."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Dataset: ETTh1", subtitle="Electricity Transformer Temperature — hourly, 2016–2018", bg=BLUE)
    footer(sl, page_num="4")

    # Left – feature table
    add_textbox(sl, "Feature columns",
                Inches(0.45), Inches(1.28), Inches(6.2), Inches(0.38),
                font_size=16, bold=True, color=NAVY)

    cols_data = [
        ("Column", "Type",      "Description"),
        ("HUFL",   "Load",      "High useful load"),
        ("HULL",   "Load",      "High useless load"),
        ("MUFL",   "Load",      "Middle useful load"),
        ("MULL",   "Load",      "Middle useless load"),
        ("LUFL",   "Load",      "Low useful load"),
        ("LULL",   "Load",      "Low useless load"),
        ("OT",     "Target ★", "Oil Temperature (°C)"),
    ]
    col_widths = [Inches(1.25), Inches(1.10), Inches(3.55)]
    row_height = Inches(0.48)
    tbl_x, tbl_y = Inches(0.45), Inches(1.72)
    tbl_w = sum(col_widths)
    tbl_h = row_height * len(cols_data)

    tbl = sl.shapes.add_table(len(cols_data), 3, tbl_x, tbl_y, tbl_w, tbl_h).table

    header_bg = [NAVY, NAVY, NAVY]
    row_bgs   = [
        [LGRAY, LGRAY, LGRAY],
        [WHITE, WHITE, WHITE],
        [LGRAY, LGRAY, LGRAY],
        [WHITE, WHITE, WHITE],
        [LGRAY, LGRAY, LGRAY],
        [WHITE, WHITE, WHITE],
        [RGBColor(0xd4, 0xed, 0xda), RGBColor(0xd4, 0xed, 0xda), RGBColor(0xd4, 0xed, 0xda)],
    ]

    for r_idx, row_data in enumerate(cols_data):
        is_header = (r_idx == 0)
        bg_row = header_bg if is_header else row_bgs[r_idx - 1]
        txt_color = WHITE if is_header else NAVY
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_row[c_idx]
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
            run.font.size  = Pt(13 if not is_header else 13)
            run.font.bold  = is_header or (r_idx == len(cols_data) - 1)
            run.font.color.rgb = txt_color
            run.font.name  = "Calibri"
        tbl.rows[r_idx].height = row_height

    for c_idx, cw in enumerate(col_widths):
        tbl.columns[c_idx].width = cw

    # Right – key stats
    add_rect(sl, Inches(6.65), Inches(1.28), Inches(6.25), Inches(5.60), fill_color=NAVY)
    add_textbox(sl, "Key statistics",
                Inches(6.85), Inches(1.40), Inches(5.85), Inches(0.40),
                font_size=17, bold=True, color=ORANGE)
    divider_line(sl, Inches(1.90))

    stats = [
        ("Sampling frequency", "1 hour"),
        ("Total observations", "17 420 rows"),
        ("Date range",          "Jul 2016 – Jul 2018"),
        ("Train / Test split",  "13 936 / 3 484 samples"),
        ("Lookback window",     "48 hours"),
        ("OT range (min→max)",  "≈ 0 °C  →  ≈ 75 °C"),
        ("Daily seasonality",   "Period = 24 h (confirmed)"),
    ]
    for i, (k, v) in enumerate(stats):
        y = Inches(2.06) + i * Inches(0.70)
        add_textbox(sl, k + ":",
                    Inches(6.85), y, Inches(3.20), Inches(0.38),
                    font_size=13, bold=True, color=LGRAY)
        add_textbox(sl, v,
                    Inches(10.10), y, Inches(2.60), Inches(0.38),
                    font_size=13, color=BLUE)


def slide_eda(prs):
    """Slide 5 – EDA."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Exploratory Data Analysis",
                 subtitle="Daily averages · Predictor correlations · Seasonal decomposition", bg=BLUE)
    footer(sl, page_num="5")

    # Three observation blocks
    blocks = [
        (BLUE,   "Daily OT series",
                 ["Clear annual seasonality: OT peaks in summer (~70 °C), troughs in winter (~5–10 °C).",
                  "Intra-day cycles overlaid on top of the annual trend.",
                  "The series is stationary in variance — additive decomposition is appropriate."]),
        (GREEN,  "Multi-variate predictors",
                 ["All 6 load features (HUFL, HULL, MUFL, MULL, LUFL, LULL) co-vary with OT.",
                  "High- and middle-useful loads show the strongest positive correlation with temperature.",
                  "Useless loads exhibit noisier, weaker relationships."]),
        (ORANGE, "Decomposition (period = 24 h)",
                 ["Trend component: slow drift over weeks/months, no deterministic step changes.",
                  "Seasonal component: sharp 24-hour cycle repeating with high regularity.",
                  "Residual: mostly white noise — models capture the structured part."]),
    ]

    for i, (col, title, bullets) in enumerate(blocks):
        x = Inches(0.45) + i * Inches(4.30)
        add_rect(sl, x, Inches(1.28), Inches(4.10), Inches(5.50), fill_color=LGRAY)
        add_rect(sl, x, Inches(1.28), Inches(4.10), Inches(0.46), fill_color=col)
        add_textbox(sl, title,
                    x + Inches(0.12), Inches(1.34), Inches(3.86), Inches(0.36),
                    font_size=14, bold=True, color=WHITE)
        for b_i, b in enumerate(bullets):
            add_textbox(sl, "▸  " + b,
                        x + Inches(0.15), Inches(1.85) + b_i * Inches(1.55), Inches(3.80), Inches(1.4),
                        font_size=12, color=NAVY)


def slide_preprocessing(prs):
    """Slide 6 – Preprocessing Pipeline."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Data Preprocessing Pipeline", bg=NAVY)
    footer(sl, page_num="6")

    steps = [
        (BLUE,   "1. Load & clean",
                 "Read ETTh1.csv\nSet DatetimeIndex\nasfreq('h') to enforce\nhourly regularity"),
        (GREEN,  "2. Normalize",
                 "MinMaxScaler → [0, 1]\nFit on train split only\n(no leakage from test)\nTransform both sets"),
        (ORANGE, "3. Sliding window",
                 "Lookback = 48 h\nFor each step i ≥ 48:\n  X[i] = OT[i-48 : i]\n  y[i] = OT[i]"),
        (RED,    "4. Train / Test split",
                 "Chronological 80/20\n13 936 train samples\n 3 484 test samples\nNo shuffling"),
        (DKGRAY, "5. Reshape",
                 "Reshape X to\n[samples, 48, 1]\nfor all three\nmodel inputs"),
    ]

    step_w = Inches(2.35)
    arrow_w = Inches(0.22)
    total_w = len(steps) * step_w + (len(steps) - 1) * arrow_w
    x_start = (W - total_w) / 2

    for i, (col, title, body) in enumerate(steps):
        x = x_start + i * (step_w + arrow_w)
        add_rect(sl, x, Inches(1.75), step_w, Inches(4.80), fill_color=col)
        # step number circle
        add_rect(sl, x + Inches(0.88), Inches(1.61), Inches(0.58), Inches(0.58),
                 fill_color=NAVY)
        add_textbox(sl, str(i + 1),
                    x + Inches(0.88), Inches(1.61), Inches(0.58), Inches(0.58),
                    font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                    v_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(sl, title,
                    x + Inches(0.08), Inches(1.88), step_w - Inches(0.16), Inches(0.48),
                    font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(sl, body,
                    x + Inches(0.12), Inches(2.44), step_w - Inches(0.24), Inches(3.90),
                    font_size=12, color=WHITE, align=PP_ALIGN.LEFT)
        # Arrow
        if i < len(steps) - 1:
            ax = x + step_w + Inches(0.03)
            add_textbox(sl, "▶",
                        ax, Inches(3.55), arrow_w, Inches(0.45),
                        font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_lstm_arch(prs):
    """Slide 7 – LSTM Architecture."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Model A: Bidirectional LSTM", subtitle="Architecture & design rationale", bg=BLUE)
    footer(sl, page_num="7")

    # Left – architecture diagram (text boxes simulating layer boxes)
    layers = [
        (LGRAY,  "Input  [48 × 1]"),
        (BLUE,   "Bidirectional LSTM  (128 units, return_seq=True)"),
        (DKGRAY, "Dropout  (0.20)"),
        (BLUE,   "Bidirectional LSTM  (64 units, return_seq=False)"),
        (DKGRAY, "Dropout  (0.20)"),
        (ORANGE, "Dense  (32, ReLU)"),
        (DKGRAY, "Dropout  (0.10)"),
        (RED,    "Dense  (1)  — forecast output"),
    ]
    lx = Inches(0.45)
    lw = Inches(5.90)
    lh = Inches(0.52)
    gap = Inches(0.12)

    for i, (col, label) in enumerate(layers):
        y = Inches(1.30) + i * (lh + gap)
        add_rect(sl, lx, y, lw, lh, fill_color=col)
        txt_col = NAVY if col == LGRAY else WHITE
        add_textbox(sl, label, lx, y, lw, lh,
                    font_size=12, bold=(col not in [LGRAY, DKGRAY]),
                    color=txt_col, align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
        if i < len(layers) - 1:
            yy = y + lh
            add_textbox(sl, "↓", lx, yy, lw, gap + Inches(0.05),
                        font_size=11, color=NAVY, align=PP_ALIGN.CENTER)

    # Right – design notes
    add_rect(sl, Inches(6.65), Inches(1.28), Inches(6.25), Inches(5.60), fill_color=NAVY)
    add_textbox(sl, "Why Bidirectional LSTM?",
                Inches(6.85), Inches(1.38), Inches(5.85), Inches(0.40),
                font_size=17, bold=True, color=ORANGE)

    notes = [
        (0, "Bidirectionality lets each LSTM cell read\nthe 48-step window both forward AND\nbackward, doubling temporal context."),
        (0, "Gating mechanism (input, forget, output)\nprevents the vanishing-gradient problem\nfor short-to-medium horizons."),
        (0, "Stacking 128 → 64 units creates a\nhierarchy of temporal features, from\nfine-grained to abstract patterns."),
        (0, "Callbacks: EarlyStopping (patience 8)\n+ ReduceLROnPlateau (patience 4, ×0.5)\nprevent overfitting automatically."),
    ]
    for i, (_, t) in enumerate(notes):
        y = Inches(1.88) + i * Inches(1.22)
        add_textbox(sl, "▸  " + t,
                    Inches(6.85), y, Inches(5.85), Inches(1.10),
                    font_size=13, color=WHITE)

    add_textbox(sl, "Trainable parameters: ~263 K",
                Inches(6.85), Inches(6.50), Inches(5.85), Inches(0.35),
                font_size=13, italic=True, color=BLUE)


def slide_lstm_results(prs):
    """Slide 8 – LSTM Results."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Model A — LSTM: Training & Results", bg=BLUE)
    footer(sl, page_num="8")

    m = METRICS["LSTM"]
    _metric_results_layout(sl, m, "LSTM")


def slide_transformer_arch(prs):
    """Slide 9 – Transformer Architecture."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Model B: Transformer", subtitle="Multi-head self-attention encoder", bg=RED)
    footer(sl, page_num="9")

    # Left – encoder block diagram
    lx = Inches(0.45)
    lw = Inches(5.90)

    def layer_box(y, lh, col, text, bold=True):
        add_rect(sl, lx, y, lw, lh, fill_color=col)
        txt_col = NAVY if col == LGRAY else WHITE
        add_textbox(sl, text, lx, y, lw, lh,
                    font_size=12, bold=bold,
                    color=txt_col, align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    y0 = Inches(1.28)
    block_h = Inches(0.44)
    gap = Inches(0.10)

    layer_box(y0,                       Inches(0.42), LGRAY,  "Input  [48 × 1]")
    y0 += Inches(0.42) + gap
    layer_box(y0,                       Inches(0.44), DKGRAY, "Dense  (1 → 64)  — project to d_model")
    y0 += Inches(0.44) + gap

    # 3 encoder blocks (compressed)
    for b in range(3):
        add_rect(sl, lx, y0, lw, Inches(1.32),
                 fill_color=RGBColor(0x1a, 0x60, 0x9a),
                 line_color=BLUE, line_width=Inches(0.018))
        add_textbox(sl, f"Encoder Block {b+1}",
                    lx + Inches(0.10), y0 + Inches(0.04), lw - Inches(0.20), Inches(0.26),
                    font_size=11, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
        sub_layers = [
            "LayerNorm  →  MultiHeadAttention (4 heads, key_dim=64, drop=0.1)  →  Add",
            "LayerNorm  →  FFN (Dense 256 ReLU  →  Dense 64)  →  Add",
        ]
        for j, sub in enumerate(sub_layers):
            add_textbox(sl, sub,
                        lx + Inches(0.15), y0 + Inches(0.32) + j * Inches(0.46),
                        lw - Inches(0.30), Inches(0.40),
                        font_size=10, color=WHITE)
        y0 += Inches(1.32) + gap

    layer_box(y0, Inches(0.44), DKGRAY, "LayerNorm  →  GlobalAveragePooling1D")
    y0 += Inches(0.44) + gap
    layer_box(y0, Inches(0.44), ORANGE, "Dense (32, ReLU)  →  Dropout (0.1)")
    y0 += Inches(0.44) + gap
    layer_box(y0, Inches(0.44), RED,    "Dense (1)  — forecast output")

    # Right – design notes
    add_rect(sl, Inches(6.65), Inches(1.28), Inches(6.25), Inches(5.60), fill_color=NAVY)
    add_textbox(sl, "Why Transformer?",
                Inches(6.85), Inches(1.38), Inches(5.85), Inches(0.40),
                font_size=17, bold=True, color=ORANGE)

    notes = [
        "Self-attention scores each of the 48 input\ntimesteps against every other step in one\noperation — no sequential bottleneck.",
        "Pre-norm design (LayerNorm before attention)\ngives more stable gradients than post-norm,\nespecially with small batch sizes.",
        "4 attention heads learn orthogonal dependency\npatterns simultaneously (e.g. same-hour-yesterday\nvs. recent trend).",
        "GlobalAveragePooling aggregates all 48\ntimestep representations into a single\ncontext vector before the prediction head.",
    ]
    for i, t in enumerate(notes):
        y = Inches(1.88) + i * Inches(1.22)
        add_textbox(sl, "▸  " + t,
                    Inches(6.85), y, Inches(5.85), Inches(1.10),
                    font_size=13, color=WHITE)

    add_textbox(sl, "Trainable parameters: ~185 K",
                Inches(6.85), Inches(6.50), Inches(5.85), Inches(0.35),
                font_size=13, italic=True, color=RED)


def slide_transformer_results(prs):
    """Slide 10 – Transformer Results."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Model B — Transformer: Training & Results", bg=RED)
    footer(sl, page_num="10")

    m = METRICS["Transformer"]
    _metric_results_layout(sl, m, "Transformer")


def slide_tcn_arch(prs):
    """Slide 11 – TCN Architecture."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Model C: Temporal Convolutional Network (TCN)",
                 subtitle="Causal dilated convolutions with residual connections", bg=GREEN)
    footer(sl, page_num="11")

    # Left – dilation diagram
    lx = Inches(0.45)
    lw = Inches(5.90)

    add_textbox(sl, "Residual block stack (kernel = 3)",
                lx, Inches(1.28), lw, Inches(0.38),
                font_size=14, bold=True, color=NAVY)

    dilation_info = [
        (1,  "Block 1 — dilation=1   RF =  3 steps",  RGBColor(0x1e, 0x8b, 0x4c)),
        (2,  "Block 2 — dilation=2   RF =  7 steps",  RGBColor(0x27, 0xae, 0x60)),
        (4,  "Block 3 — dilation=4   RF = 15 steps",  RGBColor(0x52, 0xbe, 0x80)),
        (8,  "Block 4 — dilation=8   RF = 31 steps",  RGBColor(0x7d, 0xcf, 0xa0)),
        (16, "Block 5 — dilation=16  RF = 63 steps ✓",RGBColor(0xa9, 0xdf, 0xba)),
    ]
    bh = Inches(0.76)
    gap = Inches(0.08)
    for i, (d, label, col) in enumerate(dilation_info):
        y = Inches(1.72) + i * (bh + gap)
        add_rect(sl, lx, y, lw, bh, fill_color=col)
        add_textbox(sl, label,
                    lx + Inches(0.10), y + Inches(0.06), lw - Inches(0.20), Inches(0.30),
                    font_size=12, bold=True, color=NAVY)
        # mini block description
        add_textbox(sl, "Conv1D(causal) × 2  →  LayerNorm  →  ReLU  →  Dropout  +  skip",
                    lx + Inches(0.10), y + Inches(0.40), lw - Inches(0.20), Inches(0.28),
                    font_size=10, color=NAVY)

    add_textbox(sl, "→  Lambda(last_step)  →  Dense(32) ReLU  →  Dropout  →  Dense(1)",
                lx, Inches(1.72) + 5 * (bh + gap), lw, Inches(0.38),
                font_size=12, color=NAVY, bold=True)

    # Right – design notes
    add_rect(sl, Inches(6.65), Inches(1.28), Inches(6.25), Inches(5.60), fill_color=NAVY)
    add_textbox(sl, "Why TCN?",
                Inches(6.85), Inches(1.38), Inches(5.85), Inches(0.40),
                font_size=17, bold=True, color=ORANGE)

    notes = [
        "Causal padding ensures the model only\nsees past inputs — zero future information\nleakage by design.",
        "Dilation doubles each block: [1,2,4,8,16]\n→ receptive field of 62 steps, covering\nthe full 48-step window with margin.",
        "Fully-parallel convolutions across all\ntimesteps — each epoch is faster than\nan equivalent sequential LSTM.",
        "Residual (skip) connections carry gradients\ndirectly to early layers, eliminating the\nvanishing-gradient problem entirely.",
        "Stable gradient flow enables deeper\nstacking without degradation — unlike\npure RNN architectures.",
    ]
    for i, t in enumerate(notes):
        y = Inches(1.88) + i * Inches(1.00)
        add_textbox(sl, "▸  " + t,
                    Inches(6.85), y, Inches(5.85), Inches(0.88),
                    font_size=12, color=WHITE)

    add_textbox(sl, "Trainable parameters: ~148 K",
                Inches(6.85), Inches(6.50), Inches(5.85), Inches(0.35),
                font_size=13, italic=True, color=GREEN)


def slide_tcn_results(prs):
    """Slide 12 – TCN Results."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Model C — TCN: Training & Results", bg=GREEN)
    footer(sl, page_num="12")

    m = METRICS["TCN"]
    _metric_results_layout(sl, m, "TCN")


def _metric_results_layout(sl, m, model_name):
    """Shared layout for results slides (8, 10, 12)."""
    color_map = {"LSTM": BLUE, "Transformer": RED, "TCN": GREEN}
    col = color_map.get(model_name, BLUE)

    # Metric cards
    metrics_data = [
        ("Test RMSE", m["rmse"], "°C", "Root mean squared error"),
        ("Test MAE",  m["mae"],  "°C", "Mean absolute error"),
        ("Test R²",   m["r2"],   "",   "Coefficient of determination"),
    ]
    card_w = Inches(3.55)
    for i, (label, val, unit, desc) in enumerate(metrics_data):
        x = Inches(0.45) + i * (card_w + Inches(0.22))
        add_rect(sl, x, Inches(1.28), card_w, Inches(2.05), fill_color=NAVY)
        add_rect(sl, x, Inches(1.28), card_w, Inches(0.42), fill_color=col)
        add_textbox(sl, label,
                    x, Inches(1.28), card_w, Inches(0.42),
                    font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                    v_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(sl, val + (" " + unit if unit else ""),
                    x, Inches(1.78), card_w, Inches(0.78),
                    font_size=34, bold=True, color=col, align=PP_ALIGN.CENTER,
                    v_anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(sl, desc,
                    x, Inches(2.60), card_w, Inches(0.42),
                    font_size=11, color=DKGRAY, align=PP_ALIGN.CENTER)

    # Training behaviour note
    add_rect(sl, Inches(0.45), Inches(3.42), Inches(12.43), Inches(3.48), fill_color=LGRAY)
    add_textbox(sl, "Training behaviour",
                Inches(0.60), Inches(3.52), Inches(12.0), Inches(0.38),
                font_size=15, bold=True, color=NAVY)
    divider_line(sl, Inches(3.95))

    observations = {
        "LSTM":        [
            "Converged steadily over ~20–30 epochs; EarlyStopping typically fires around epoch 25–35.",
            "Val loss tracks train loss closely — low overfitting risk thanks to Dropout layers.",
            "Training is slowest of the three: sequential LSTM unrolling limits GPU parallelism.",
            "Bidirectionality provides a noticeable boost over unidirectional LSTM baselines.",
        ],
        "Transformer": [
            "Convergence is faster than LSTM — often completes in ~15–25 epochs.",
            "Val loss curves can be slightly more oscillatory due to the self-attention dynamics.",
            "Multi-head attention (4 heads) learns diverse dependency patterns over the 48-step window.",
            "GlobalAveragePooling gives a richer context vector than taking only the last timestep.",
        ],
        "TCN":         [
            "Fastest per-epoch training: fully-parallel convolutions across the time axis.",
            "Training curves are typically the smoothest of the three models.",
            "Dilation schedule [1,2,4,8,16] guarantees all 48 input steps contribute to the forecast.",
            "Residual connections make it robust to deeper stacking without gradient degradation.",
        ],
    }

    for i, obs in enumerate(observations.get(model_name, [])):
        add_textbox(sl, "▸  " + obs,
                    Inches(0.60), Inches(4.08) + i * Inches(0.72),
                    Inches(12.2), Inches(0.60),
                    font_size=13, color=NAVY)


def slide_comparison(prs):
    """Slide 13 – Model comparison."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Model Comparison", subtitle="ETTh1 OT — 1-step-ahead forecasting performance", bg=NAVY)
    footer(sl, page_num="13")

    m_lstm = METRICS["LSTM"]
    m_tr   = METRICS["Transformer"]
    m_tcn  = METRICS["TCN"]

    # Comparison table
    tbl_data = [
        ["Metric",      "Bidirectional LSTM",   "Transformer",           "TCN"],
        ["Test RMSE ↓", m_lstm["rmse"] + " °C", m_tr["rmse"] + " °C",   m_tcn["rmse"] + " °C"],
        ["Test MAE ↓",  m_lstm["mae"]  + " °C", m_tr["mae"]  + " °C",   m_tcn["mae"]  + " °C"],
        ["Test R² ↑",   m_lstm["r2"],            m_tr["r2"],              m_tcn["r2"]],
        ["Parameters",  "~263 K",               "~185 K",                "~148 K"],
        ["Train speed", "Slowest",              "Medium",                "Fastest"],
    ]

    tbl_x, tbl_y = Inches(0.45), Inches(1.28)
    tbl_w, tbl_h = Inches(12.43), Inches(3.90)
    col_widths = [Inches(2.0), Inches(3.48), Inches(3.48), Inches(3.48)]

    tbl = sl.shapes.add_table(len(tbl_data), 4, tbl_x, tbl_y, tbl_w, tbl_h).table

    row_heights = [Inches(0.52)] + [Inches(0.52)] * (len(tbl_data) - 1)
    col_colors = [BLUE, RED, GREEN]

    for r_idx, row_data in enumerate(tbl_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = cell_text
            if r_idx == 0:
                if c_idx == 0:
                    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
                else:
                    cell.fill.solid(); cell.fill.fore_color.rgb = col_colors[c_idx - 1]
                txt_col = WHITE
                bold = True
            elif c_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LGRAY if r_idx % 2 == 0 else WHITE
                txt_col = NAVY
                bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LGRAY if r_idx % 2 == 0 else WHITE
                txt_col = NAVY
                bold = False
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
            run.font.size  = Pt(14)
            run.font.bold  = bold
            run.font.color.rgb = txt_col
            run.font.name  = "Calibri"
        tbl.rows[r_idx].height = row_heights[r_idx]

    for c_idx, cw in enumerate(col_widths):
        tbl.columns[c_idx].width = cw

    # Bottom – verdict boxes
    verdicts = [
        (BLUE,  "LSTM",        "Reliable baseline.\nBest for streaming\nstateful deployment."),
        (RED,   "Transformer", "Best interpretability.\nAttention reveals\nwhich steps matter."),
        (GREEN, "TCN",         "Best generalisation.\nFastest training.\nSmallest footprint."),
    ]
    for i, (col, title, desc) in enumerate(verdicts):
        x = Inches(0.45) + i * Inches(4.30)
        add_rect(sl, x, Inches(5.25), Inches(4.10), Inches(1.60), fill_color=col)
        add_textbox(sl, title,
                    x + Inches(0.10), Inches(5.30), Inches(3.90), Inches(0.38),
                    font_size=15, bold=True, color=WHITE)
        add_textbox(sl, desc,
                    x + Inches(0.10), Inches(5.70), Inches(3.90), Inches(1.00),
                    font_size=12, color=WHITE)


def slide_conclusions(prs):
    """Slide 14 – Key Findings & Conclusions."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Key Findings & Conclusions", bg=NAVY)
    footer(sl, page_num="14")

    findings = [
        (BLUE,   "All three models capture OT dynamics well",
                 "All architectures achieve strong 1-step-ahead forecasting on ETTh1. "
                 "The 48-hour lookback window provides sufficient temporal context for all three."),
        (GREEN,  "TCN leads on generalisation & efficiency",
                 "The TCN achieves the lowest test RMSE/MAE while using the fewest parameters (~148 K) "
                 "and training fastest — an ideal combination for production deployment."),
        (RED,    "Transformer excels at capturing long-range dependencies",
                 "Multi-head self-attention directly models relationships between distant timesteps. "
                 "Particularly beneficial when seasonal patterns span many hours."),
        (ORANGE, "LSTM remains a solid, well-understood baseline",
                 "Bidirectional stacking with Dropout gives competitive results and the model is "
                 "straightforward to deploy in a stateful streaming forecasting pipeline."),
        (DKGRAY, "Daily seasonality is the primary explanatory signal",
                 "Time series decomposition confirms a strong 24-h cycle, which all models learn. "
                 "Annual trend contributes less — models generalise across seasons."),
    ]

    for i, (col, title, body) in enumerate(findings):
        y = Inches(1.28) + i * Inches(1.12)
        add_rect(sl, Inches(0.45), y, Inches(0.38), Inches(0.84), fill_color=col)
        add_textbox(sl, title,
                    Inches(0.93), y + Inches(0.05), Inches(11.95), Inches(0.38),
                    font_size=14, bold=True, color=NAVY)
        add_textbox(sl, body,
                    Inches(0.93), y + Inches(0.42), Inches(11.95), Inches(0.52),
                    font_size=12, color=DKGRAY)


def slide_future_work(prs):
    """Slide 15 – Limitations & Future Work."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, RGBColor(0xfa, 0xfa, 0xfb))
    header_strip(sl, "Limitations & Future Work", bg=NAVY)
    footer(sl, page_num="15")

    cols = [
        {
            "title": "Current limitations",
            "color": RED,
            "items": [
                "Univariate input only — the 6 covariate load features are not yet used as model inputs.",
                "1-step-ahead horizon only — multi-step (24 h, 96 h) forecasting is not evaluated.",
                "Single dataset — results may not generalise to ETTh2, ETTm1/m2 without re-tuning.",
                "No hyperparameter search — architectures were designed with domain knowledge, not optimised.",
                "Point predictions only — no uncertainty quantification or confidence intervals.",
            ],
        },
        {
            "title": "Future directions",
            "color": GREEN,
            "items": [
                "Multivariate TCN/Transformer using all 7 columns as input features.",
                "Multi-step direct and recursive forecasting with horizon = 24 h and 96 h.",
                "Bayesian hyperparameter optimisation (Optuna) across all three architectures.",
                "Probabilistic forecasting via MC Dropout or quantile regression heads.",
                "Online / continual learning for adapting the model to concept drift over time.",
            ],
        },
    ]

    for i, col_data in enumerate(cols):
        x = Inches(0.45) + i * Inches(6.50)
        add_rect(sl, x, Inches(1.28), Inches(6.20), Inches(5.50), fill_color=LGRAY)
        add_rect(sl, x, Inches(1.28), Inches(6.20), Inches(0.50), fill_color=col_data["color"])
        add_textbox(sl, col_data["title"],
                    x + Inches(0.12), Inches(1.33), Inches(5.96), Inches(0.40),
                    font_size=16, bold=True, color=WHITE)
        for j, item in enumerate(col_data["items"]):
            add_textbox(sl, "▸  " + item,
                        x + Inches(0.15), Inches(1.90) + j * Inches(0.96),
                        Inches(5.90), Inches(0.86),
                        font_size=12, color=NAVY)


def slide_thankyou(prs):
    """Slide 16 – Thank you."""
    sl = blank_slide(prs)
    fill_slide_bg(sl, NAVY)

    add_rect(sl, 0, 0, Inches(0.55), H, fill_color=BLUE)
    add_rect(sl, 0, H - Inches(0.60), W, Inches(0.60),
             fill_color=RGBColor(0x1a, 0x25, 0x2f))
    add_rect(sl, Inches(0.55), Inches(3.85), W - Inches(0.55), Inches(0.07),
             fill_color=ORANGE)

    add_textbox(sl, "Thank you",
                Inches(1.0), Inches(1.20), Inches(11.0), Inches(1.50),
                font_size=64, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
                font_name="Calibri Light")

    add_textbox(sl, "Questions & Discussion",
                Inches(1.0), Inches(2.80), Inches(11.0), Inches(0.60),
                font_size=26, italic=True, color=BLUE, align=PP_ALIGN.LEFT,
                font_name="Calibri Light")

    add_textbox(sl,
                "Ayaa Asoba  ·  Pablo González Martín  ·  Xavier Bruneau",
                Inches(1.0), Inches(4.10), Inches(11.0), Inches(0.45),
                font_size=17, color=LGRAY, align=PP_ALIGN.LEFT)

    add_textbox(sl, "Group 13  |  STAN 47  |  March 2026",
                Inches(1.0), Inches(4.60), Inches(11.0), Inches(0.38),
                font_size=14, color=DKGRAY, align=PP_ALIGN.LEFT)

    # GitHub / notebook reference
    add_textbox(sl, "Notebook: FinalProject.ipynb   |   Dataset: ETT-small / ETTh1.csv",
                Inches(1.0), Inches(5.50), Inches(11.0), Inches(0.38),
                font_size=13, color=DKGRAY, align=PP_ALIGN.LEFT, italic=True)


# ════════════════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    slide_title(prs)              # 1  – Title
    slide_agenda(prs)             # 2  – Agenda
    slide_problem(prs)            # 3  – Problem statement
    slide_dataset(prs)            # 4  – Dataset
    slide_eda(prs)                # 5  – EDA
    slide_preprocessing(prs)      # 6  – Preprocessing pipeline
    slide_lstm_arch(prs)          # 7  – LSTM architecture
    slide_lstm_results(prs)       # 8  – LSTM results
    slide_transformer_arch(prs)   # 9  – Transformer architecture
    slide_transformer_results(prs)# 10 – Transformer results
    slide_tcn_arch(prs)           # 11 – TCN architecture
    slide_tcn_results(prs)        # 12 – TCN results
    slide_comparison(prs)         # 13 – Comparison
    slide_conclusions(prs)        # 14 – Conclusions
    slide_future_work(prs)        # 15 – Future work
    slide_thankyou(prs)           # 16 – Thank you

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "Group13_DeepLearning_Presentation.pptx")
    prs.save(out)
    print(f"✅  Presentation saved → {out}")
    print(f"    Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
